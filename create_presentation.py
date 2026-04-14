#!/usr/bin/env python3
"""Generate PowerPoint presentation for EduFeedback project"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
DARK_BLUE = RGBColor(26, 26, 46)  # #1a1a2e
INDIGO = RGBColor(99, 102, 241)   # #6366f1
EMERALD = RGBColor(16, 185, 129)  # #10b981
WHITE = RGBColor(241, 245, 249)   # #f1f5f9
GRAY = RGBColor(100, 116, 139)    # #64748b

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(66)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.text = subtitle
    subtitle_p.font.size = Pt(24)
    subtitle_p.font.color.rgb = INDIGO
    
    return slide

def add_content_slide(prs, title, content_points):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = DARK_BLUE
    title_shape.line.color.rgb = DARK_BLUE
    
    # Title text
    title_frame = title_shape.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(44)
    title_p.font.bold = True
    title_p.font.color.rgb = INDIGO
    title_p.space_before = Pt(12)
    title_p.space_after = Pt(12)
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.4), Inches(5.8))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, point in enumerate(content_points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = point
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_BLUE
        p.space_before = Pt(10)
        p.space_after = Pt(10)
        p.level = 0
        
    return slide

# Slide 1: Title
add_title_slide(prs, "🎓 EduFeedback", "Course Feedback & Rating System\nA Full-Stack MERN Application")

# Slide 2: Project Overview
add_content_slide(prs, "Project Overview", [
    "📌 Replaces paper-based course feedback with a modern digital platform",
    "🔐 Secure authentication and user management system",
    "⭐ Star-based rating system for courses",
    "👥 Role-based access control (Student, Admin)",
    "📊 Real-time feedback analytics and visualization",
    "🚀 Built with MERN Stack (MongoDB, Express, React, Node.js)"
])

# Slide 3: Key Features
add_content_slide(prs, "Key Features", [
    "✅ User Authentication & Registration with JWT",
    "✅ Student Dashboard with course listing",
    "✅ Star Rating System for course feedback",
    "✅ Admin Panel for analytics and course management",
    "✅ Secure API with middleware protection",
    "✅ Real-time notifications with React Hot Toast"
])

# Slide 4: Architecture
add_content_slide(prs, "System Architecture", [
    "Frontend: React.js with Context API for state management",
    "Backend: Express.js REST API with middleware",
    "Database: MongoDB with Mongoose ODM",
    "Authentication: JWT tokens with secure storage",
    "API Routes: Auth, Courses, Feedback management"
])

# Slide 5: Tech Stack - Frontend
add_content_slide(prs, "Tech Stack - Frontend", [
    "⚡ React.js - UI framework",
    "🛣️ React Router - Client-side routing",
    "🎨 CSS3 - Modern styling",
    "🔔 React Hot Toast - Notifications",
    "🌐 Axios - HTTP client",
    "📦 Context API - State management"
])

# Slide 6: Tech Stack - Backend
add_content_slide(prs, "Tech Stack - Backend", [
    "✨ Node.js - Runtime environment",
    "🚀 Express.js - Web framework",
    "🗄️ MongoDB - NoSQL database",
    "🔐 JWT - Authentication tokens",
    "🛡️ CORS - Cross-origin security",
    "📝 Mongoose - Schema validation"
])

# Slide 7: Database Models
add_content_slide(prs, "Database Models", [
    "👤 User Model: ID, name, email, password, role (student/admin)",
    "📚 Course Model: ID, title, instructor, semester, code",
    "⭐ Feedback Model: ID, userId, courseId, rating (1-5), comments, timestamp",
    "🔗 Relationships: Feedback references User and Course"
])

# Slide 8: User Roles & Permissions
add_content_slide(prs, "User Roles & Permissions", [
    "👨‍🎓 STUDENT: View courses, submit feedback, view ratings",
    "👨‍💼 ADMIN: View analytics, manage courses, approve feedbacks",
    "🔐 Protected Routes: Middleware-based access control",
    "📋 Route Guards: ProtectedRoute & AdminRoute components"
])

# Slide 9: Project File Structure
add_content_slide(prs, "Project Structure", [
    "server/ → Express backend with routes, controllers, models",
    "client/ → React frontend with components and services",
    "server/config/ → Database configuration",
    "server/middleware/ → Auth & admin middleware",
    "client/pages/ → Login, Dashboard, Feedback, Admin Panel",
    "client/services/ → API service modules"
])

# Slide 10: Installation & Setup
add_content_slide(prs, "Installation & Setup", [
    "1️⃣ Install Node.js v16+ and MongoDB",
    "2️⃣ Backend: cd server && npm install",
    "3️⃣ Frontend: cd client && npm install",
    "4️⃣ Configure environment variables",
    "5️⃣ Run: npm start (backend) && npm start (frontend)",
    "6️⃣ Access at http://localhost:3000"
])

# Slide 11: API Endpoints
add_content_slide(prs, "API Endpoints Overview", [
    "🔐 Auth: POST /api/auth/register, POST /api/auth/login",
    "📚 Courses: GET /api/courses, POST /api/courses (admin)",
    "⭐ Feedback: POST /api/feedback/submit, GET /api/feedback/analytics",
    "🔒 All endpoints protected with JWT middleware",
    "📊 Admin endpoints require admin middleware verification"
])

# Slide 12: Demo Workflow
add_content_slide(prs, "User Workflow", [
    "1. User registers with email and password",
    "2. Login and receive JWT token",
    "3. Dashboard displays available courses",
    "4. Submit star ratings and feedback for each course",
    "5. Admin views analytics and feedback trends",
    "6. System stores all feedback securely in MongoDB"
])

# Slide 13: Security Features
add_content_slide(prs, "Security Features", [
    "🔐 JWT Token-based authentication",
    "🛡️ Password hashing with industry standards",
    "🔒 CORS enabled for secure cross-origin requests",
    "👮 Role-based middleware protection",
    "🚫 Input validation on all endpoints",
    "📝 Secure credential management via environment variables"
])

# Slide 14: Future Enhancements
add_content_slide(prs, "Future Enhancements", [
    "📱 Mobile app for iOS/Android",
    "📊 Advanced analytics and graphical reports",
    "🔔 Email notifications for feedback",
    "🌍 Multi-language support",
    "🔍 Full-text search functionality",
    "💾 Data export features (PDF, Excel)"
])

# Slide 15: Team & Credits
add_content_slide(prs, "Team & Credits", [
    "👨‍💻 Developer: A. Yashaswini (B23CN111)",
    "👨‍🏫 Supervisor: BRIKIENLABS",
    "🔗 Website: https://brikienlabs.tech",
    "📧 Support & Contributions Welcome",
    "📄 Open Source Project",
    "✨ Thank You!"
])

# Save presentation
output_path = r'c:\survey\EduFeedback_Presentation.pptx'
prs.save(output_path)
print(f"✅ Presentation created successfully: {output_path}")
